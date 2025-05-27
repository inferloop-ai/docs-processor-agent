"""
Document Loader Agent
Multi-format document loading with intelligent processing
"""

import asyncio
import hashlib
import mimetypes
import os
import tempfile
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, Tuple
import magic
from datetime import datetime

# Document processing imports
from langchain.schema import Document
from langchain.document_loaders import (
    PyPDFLoader, UnstructuredWordDocumentLoader, TextLoader,
    UnstructuredHTMLLoader, UnstructuredMarkdownLoader,
    UnstructuredPowerPointLoader, UnstructuredExcelLoader,
    CSVLoader, JSONLoader
)
from langchain.document_loaders.blob_loaders import Blob
from langchain.document_loaders.parsers import (
    OpenAIWhisperParser, PyMuPDFParser
)

# OCR and image processing
try:
    import pytesseract
    from PIL import Image
    import fitz  # PyMuPDF
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

# Office document processing
try:
    import docx2txt
    import pptx
    import openpyxl
    OFFICE_AVAILABLE = True
except ImportError:
    OFFICE_AVAILABLE = False

from .base_agent import BaseAgent, AgentResult


class DocumentLoader(BaseAgent):
    """
    Document Loader Agent
    
    Capabilities:
    - Multi-format document loading (PDF, Word, Excel, PowerPoint, Text, HTML, etc.)
    - OCR for scanned documents and images
    - Content extraction and preprocessing
    - Metadata extraction
    - File validation and security checks
    - Intelligent format detection
    - Content deduplication
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__(config)
        
        # Load document loader specific configuration
        self.loader_config = self.config.get('agents', {}).get('document_loader', {})
        
        # Supported formats
        self.supported_formats = self.loader_config.get('supported_formats', [
            '.pdf', '.docx', '.doc', '.txt', '.md', '.html', '.htm',
            '.pptx', '.ppt', '.xlsx', '.xls', '.csv', '.json', '.rtf'
        ])
        
        # File size limits
        self.max_file_size_mb = self.loader_config.get('max_file_size_mb', 100)
        self.max_file_size_bytes = self.max_file_size_mb * 1024 * 1024
        
        # Processing options
        self.use_ocr = self.loader_config.get('use_ocr', True)
        self.extract_images = self.loader_config.get('extract_images', False)
        self.preserve_formatting = self.loader_config.get('preserve_formatting', True)
        self.temp_dir = self.loader_config.get('temp_dir', tempfile.gettempdir())
        
        # Content validation
        self.min_content_length = self.loader_config.get('min_content_length', 10)
        self.max_content_length = self.loader_config.get('max_content_length', 10_000_000)
        
        # OCR settings
        self.ocr_languages = self.loader_config.get('ocr_languages', ['eng'])
        self.ocr_config = self.loader_config.get('ocr_config', '--psm 3')
        
        # Initialize magic for file type detection
        try:
            self.magic = magic.Magic(mime=True)
        except:
            self.magic = None
            self.log_warning("python-magic not available, using mimetypes fallback")
        
        self.log_info(f"DocumentLoader initialized with {len(self.supported_formats)} supported formats")
    
    async def process(self, input_data: Union[str, Path, Dict[str, Any]], **kwargs) -> AgentResult:
        """
        Load and process documents
        
        Args:
            input_data: File path, Path object, or dict with file info
            **kwargs: Additional parameters
                - file_content: Raw file content (bytes)
                - filename: Original filename
                - extract_metadata: Whether to extract detailed metadata
                - use_ocr: Override OCR setting
        
        Returns:
            AgentResult with loaded documents and metadata
        """
        try:
            # Parse input
            if isinstance(input_data, (str, Path)):
                file_path = Path(input_data)
                if not file_path.exists():
                    return AgentResult(
                        success=False,
                        error=f"File not found: {file_path}"
                    )
                
                filename = file_path.name
                file_content = None
                
            elif isinstance(input_data, dict):
                file_path = Path(input_data.get('file_path', ''))
                filename = input_data.get('filename', file_path.name if file_path else 'unknown')
                file_content = input_data.get('file_content')
                
                if not file_path.exists() and not file_content:
                    return AgentResult(
                        success=False,
                        error="Either file_path or file_content must be provided"
                    )
            else:
                return AgentResult(
                    success=False,
                    error="Invalid input: expected file path or dict with file info"
                )
            
            # Validate file
            validation_result = await self._validate_file(file_path, file_content, filename)
            if not validation_result['valid']:
                return AgentResult(
                    success=False,
                    error=validation_result['error']
                )
            
            # Detect file format
            file_format = await self._detect_format(file_path, file_content, filename)
            
            # Check if format is supported
            if file_format not in self.supported_formats:
                return AgentResult(
                    success=False,
                    error=f"Unsupported format: {file_format}. Supported: {', '.join(self.supported_formats)}"
                )
            
            # Load document content
            documents = await self._load_document_by_format(
                file_path, file_content, filename, file_format, kwargs
            )
            
            # Extract metadata
            metadata = await self._extract_file_metadata(
                file_path, file_content, filename, file_format
            )
            
            # Post-process documents
            processed_documents = await self._post_process_documents(documents, metadata)
            
            # Calculate content hash for deduplication
            content_hash = self._calculate_content_hash(processed_documents)
            
            result_data = {
                "documents": processed_documents,
                "metadata": {
                    **metadata,
                    "content_hash": content_hash,
                    "processing_timestamp": datetime.now().isoformat(),
                    "total_pages": len(processed_documents),
                    "total_characters": sum(len(doc.page_content) for doc in processed_documents)
                },
                "file_info": {
                    "filename": filename,
                    "format": file_format,
                    "size_bytes": validation_result.get('size_bytes', 0),
                    "mime_type": validation_result.get('mime_type')
                }
            }
            
            return AgentResult(
                success=True,
                data=result_data,
                confidence=0.95,
                metadata={
                    "loader_used": self._get_loader_name(file_format),
                    "ocr_used": kwargs.get('use_ocr', self.use_ocr) and self._requires_ocr(file_format),
                    "pages_processed": len(processed_documents)
                }
            )
            
        except Exception as e:
            self.log_error(f"Error loading document: {str(e)}")
            return AgentResult(
                success=False,
                error=f"Document loading failed: {str(e)}"
            )
    
    async def _validate_file(self, file_path: Path, file_content: Optional[bytes], 
                           filename: str) -> Dict[str, Any]:
        """Validate file before processing"""
        try:
            # Get file size
            if file_content:
                size_bytes = len(file_content)
            elif file_path.exists():
                size_bytes = file_path.stat().st_size
            else:
                return {"valid": False, "error": "File not accessible"}
            
            # Check file size
            if size_bytes > self.max_file_size_bytes:
                return {
                    "valid": False,
                    "error": f"File too large: {size_bytes / 1024 / 1024:.1f}MB > {self.max_file_size_mb}MB"
                }
            
            if size_bytes == 0:
                return {"valid": False, "error": "Empty file"}
            
            # Detect MIME type
            mime_type = None
            if self.magic:
                if file_content:
                    mime_type = self.magic.from_buffer(file_content)
                elif file_path.exists():
                    mime_type = self.magic.from_file(str(file_path))
            
            if not mime_type:
                mime_type, _ = mimetypes.guess_type(filename)
            
            # Basic security checks
            suspicious_extensions = ['.exe', '.bat', '.cmd', '.scr', '.vbs', '.js']
            file_ext = Path(filename).suffix.lower()
            
            if file_ext in suspicious_extensions:
                return {"valid": False, "error": f"Potentially unsafe file type: {file_ext}"}
            
            return {
                "valid": True,
                "size_bytes": size_bytes,
                "mime_type": mime_type
            }
            
        except Exception as e:
            return {"valid": False, "error": f"File validation error: {str(e)}"}
    
    async def _detect_format(self, file_path: Path, file_content: Optional[bytes], 
                           filename: str) -> str:
        """Detect document format"""
        # Try extension first
        file_ext = Path(filename).suffix.lower()
        if file_ext in self.supported_formats:
            return file_ext
        
        # Try MIME type detection
        if self.magic:
            try:
                if file_content:
                    mime_type = self.magic.from_buffer(file_content)
                elif file_path.exists():
                    mime_type = self.magic.from_file(str(file_path))
                else:
                    mime_type = None
                
                # Map MIME types to extensions
                mime_to_ext = {
                    'application/pdf': '.pdf',
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                    'application/msword': '.doc',
                    'text/plain': '.txt',
                    'text/html': '.html',
                    'text/markdown': '.md',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                    'application/vnd.ms-powerpoint': '.ppt',
                    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                    'application/vnd.ms-excel': '.xls',
                    'text/csv': '.csv',
                    'application/json': '.json'
                }
                
                detected_ext = mime_to_ext.get(mime_type)
                if detected_ext and detected_ext in self.supported_formats:
                    return detected_ext
                    
            except Exception as e:
                self.log_warning(f"MIME type detection failed: {e}")
        
        # Fallback to extension
        return file_ext
    
    async def _load_document_by_format(self, file_path: Path, file_content: Optional[bytes],
                                     filename: str, file_format: str, 
                                     options: Dict[str, Any]) -> List[Document]:
        """Load document based on format"""
        
        # Prepare file path (create temp file if working with bytes)
        working_path = file_path
        temp_file = None
        
        if file_content:
            temp_file = tempfile.NamedTemporaryFile(
                suffix=file_format,
                dir=self.temp_dir,
                delete=False
            )
            temp_file.write(file_content)
            temp_file.close()
            working_path = Path(temp_file.name)
        
        try:
            documents = []
            
            if file_format == '.pdf':
                documents = await self._load_pdf(working_path, options)
            elif file_format in ['.docx', '.doc']:
                documents = await self._load_word(working_path, options)
            elif file_format in ['.xlsx', '.xls']:
                documents = await self._load_excel(working_path, options)
            elif file_format in ['.pptx', '.ppt']:
                documents = await self._load_powerpoint(working_path, options)
            elif file_format in ['.txt', '.md']:
                documents = await self._load_text(working_path, options)
            elif file_format in ['.html', '.htm']:
                documents = await self._load_html(working_path, options)
            elif file_format == '.csv':
                documents = await self._load_csv(working_path, options)
            elif file_format == '.json':
                documents = await self._load_json(working_path, options)
            else:
                # Fallback to text loader
                documents = await self._load_text(working_path, options)
            
            return documents
            
        finally:
            # Clean up temp file
            if temp_file:
                try:
                    os.unlink(temp_file.name)
                except:
                    pass
    
    async def _load_pdf(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load PDF document"""
        try:
            loader = PyPDFLoader(str(file_path))
            documents = loader.load()
            
            # Add OCR if needed and available
            if (options.get('use_ocr', self.use_ocr) and 
                TESSERACT_AVAILABLE and self._should_use_ocr(documents)):
                
                ocr_documents = await self._apply_ocr_to_pdf(file_path)
                if ocr_documents:
                    documents = ocr_documents
            
            return documents
            
        except Exception as e:
            self.log_warning(f"PyPDF loading failed, trying alternative: {e}")
            
            # Alternative: try PyMuPDF
            try:
                documents = []
                doc = fitz.open(str(file_path))
                
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    text = page.get_text()
                    
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            'page': page_num + 1,
                            'source': str(file_path),
                            'total_pages': doc.page_count
                        }
                    ))
                
                doc.close()
                return documents
                
            except Exception as e2:
                raise Exception(f"PDF loading failed with both loaders: PyPDF2: {e}, PyMuPDF: {e2}")
    
    async def _load_word(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load Word document"""
        try:
            loader = UnstructuredWordDocumentLoader(str(file_path))
            documents = loader.load()
            return documents
            
        except Exception as e:
            self.log_warning(f"Unstructured Word loader failed, trying docx2txt: {e}")
            
            # Alternative: try docx2txt for .docx files
            if OFFICE_AVAILABLE and file_path.suffix.lower() == '.docx':
                try:
                    text = docx2txt.process(str(file_path))
                    return [Document(
                        page_content=text,
                        metadata={'source': str(file_path)}
                    )]
                except Exception as e2:
                    raise Exception(f"Word loading failed: Unstructured: {e}, docx2txt: {e2}")
            else:
                raise
    
    async def _load_excel(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load Excel document"""
        try:
            loader = UnstructuredExcelLoader(str(file_path))
            documents = loader.load()
            return documents
            
        except Exception as e:
            self.log_warning(f"Unstructured Excel loader failed: {e}")
            
            # Alternative: try openpyxl
            if OFFICE_AVAILABLE:
                try:
                    from openpyxl import load_workbook
                    
                    workbook = load_workbook(str(file_path), data_only=True)
                    documents = []
                    
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        
                        # Extract text from cells
                        content_lines = []
                        for row in sheet.iter_rows(values_only=True):
                            row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                            if row_text.strip():
                                content_lines.append(row_text)
                        
                        if content_lines:
                            documents.append(Document(
                                page_content='\n'.join(content_lines),
                                metadata={
                                    'source': str(file_path),
                                    'sheet': sheet_name
                                }
                            ))
                    
                    return documents
                    
                except Exception as e2:
                    raise Exception(f"Excel loading failed: Unstructured: {e}, openpyxl: {e2}")
            else:
                raise
    
    async def _load_powerpoint(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load PowerPoint document"""
        try:
            loader = UnstructuredPowerPointLoader(str(file_path))
            documents = loader.load()
            return documents
            
        except Exception as e:
            self.log_warning(f"Unstructured PowerPoint loader failed: {e}")
            
            # Alternative: try python-pptx
            if OFFICE_AVAILABLE:
                try:
                    from pptx import Presentation
                    
                    prs = Presentation(str(file_path))
                    documents = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        content_parts = []
                        
                        for shape in slide.shapes:
                            if hasattr(shape, 'text') and shape.text.strip():
                                content_parts.append(shape.text)
                        
                        if content_parts:
                            documents.append(Document(
                                page_content='\n'.join(content_parts),
                                metadata={
                                    'source': str(file_path),
                                    'slide': slide_num
                                }
                            ))
                    
                    return documents
                    
                except Exception as e2:
                    raise Exception(f"PowerPoint loading failed: Unstructured: {e}, python-pptx: {e2}")
            else:
                raise
    
    async def _load_text(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load text document"""
        loader = TextLoader(str(file_path), encoding='utf-8')
        return loader.load()
    
    async def _load_html(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load HTML document"""
        loader = UnstructuredHTMLLoader(str(file_path))
        return loader.load()
    
    async def _load_csv(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load CSV document"""
        loader = CSVLoader(str(file_path))
        return loader.load()
    
    async def _load_json(self, file_path: Path, options: Dict[str, Any]) -> List[Document]:
        """Load JSON document"""
        loader = JSONLoader(str(file_path), jq_schema='.', text_content=False)
        return loader.load()
    
    async def _apply_ocr_to_pdf(self, file_path: Path) -> Optional[List[Document]]:
        """Apply OCR to PDF if needed"""
        if not TESSERACT_AVAILABLE:
            return None
        
        try:
            import fitz
            
            documents = []
            doc = fitz.open(str(file_path))
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # Convert page to image
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                
                # OCR the image
                image = Image.open(io.BytesIO(img_data))
                text = pytesseract.image_to_string(
                    image, 
                    lang='+'.join(self.ocr_languages),
                    config=self.ocr_config
                )
                
                if text.strip():
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            'page': page_num + 1,
                            'source': str(file_path),
                            'total_pages': doc.page_count,
                            'ocr_applied': True
                        }
                    ))
            
            doc.close()
            return documents
            
        except Exception as e:
            self.log_warning(f"OCR processing failed: {e}")
            return None
    
    def _should_use_ocr(self, documents: List[Document]) -> bool:
        """Determine if OCR should be applied based on content quality"""
        if not documents:
            return True
        
        # Check if content seems to be mostly extracted text or needs OCR
        total_text = ' '.join([doc.page_content for doc in documents])
        
        # Heuristics for low-quality extraction (indicating need for OCR)
        if len(total_text.strip()) < 50:  # Very little text extracted
            return True
        
        # Check for unusual character patterns that might indicate extraction issues
        weird_char_ratio = sum(1 for c in total_text if not c.isprintable() and c not in '\n\t') / max(len(total_text), 1)
        if weird_char_ratio > 0.1:  # More than 10% weird characters
            return True
        
        return False
    
    def _requires_ocr(self, file_format: str) -> bool:
        """Check if file format might require OCR"""
        return file_format in ['.pdf']  # Extend as needed
    
    async def _extract_file_metadata(self, file_path: Path, file_content: Optional[bytes],
                                   filename: str, file_format: str) -> Dict[str, Any]:
        """Extract file metadata"""
        metadata = {
            'filename': filename,
            'file_format': file_format,
            'loader_version': '2.0'
        }
        
        try:
            # Basic file stats
            if file_path.exists():
                stat = file_path.stat()
                metadata.update({
                    'file_size': stat.st_size,
                    'created_time': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                    'modified_time': datetime.fromtimestamp(stat.st_mtime).isoformat()
                })
            
            # Format-specific metadata extraction
            if file_format == '.pdf':
                metadata.update(await self._extract_pdf_metadata(file_path))
            elif file_format in ['.docx', '.xlsx', '.pptx']:
                metadata.update(await self._extract_office_metadata(file_path, file_format))
            
        except Exception as e:
            self.log_warning(f"Metadata extraction failed: {e}")
        
        return metadata
    
    async def _extract_pdf_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract PDF-specific metadata"""
        try:
            import fitz
            
            doc = fitz.open(str(file_path))
            metadata = doc.metadata
            
            result = {
                'page_count': doc.page_count,
                'title': metadata.get('title', ''),
                'author': metadata.get('author', ''),
                'subject': metadata.get('subject', ''),
                'creator': metadata.get('creator', ''),
                'producer': metadata.get('producer', ''),
                'creation_date': metadata.get('creationDate', ''),
                'modification_date': metadata.get('modDate', '')
            }
            
            doc.close()
            return result
            
        except Exception as e:
            self.log_warning(f"PDF metadata extraction failed: {e}")
            return {}
    
    async def _extract_office_metadata(self, file_path: Path, file_format: str) -> Dict[str, Any]:
        """Extract Office document metadata"""
        # This would require python-docx, openpyxl, python-pptx for full implementation
        return {}
    
    async def _post_process_documents(self, documents: List[Document], 
                                    metadata: Dict[str, Any]) -> List[Document]:
        """Post-process loaded documents"""
        processed_docs = []
        
        for i, doc in enumerate(documents):
            # Clean content
            cleaned_content = self._clean_text(doc.page_content)
            
            # Skip empty documents
            if len(cleaned_content.strip()) < self.min_content_length:
                continue
            
            # Truncate if too long
            if len(cleaned_content) > self.max_content_length:
                cleaned_content = cleaned_content[:self.max_content_length]
                self.log_warning(f"Document {i} truncated to {self.max_content_length} characters")
            
            # Update metadata
            enhanced_metadata = {
                **doc.metadata,
                'document_index': i,
                'character_count': len(cleaned_content),
                'word_count': len(cleaned_content.split()),
                'processing_timestamp': datetime.now().isoformat()
            }
            
            processed_docs.append(Document(
                page_content=cleaned_content,
                metadata=enhanced_metadata
            ))
        
        return processed_docs
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text content"""
        if not text:
            return ""
        
        # Basic cleaning
        text = text.strip()
        
        # Remove excessive whitespace
        import re
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)  # Multiple newlines to double
        text = re.sub(r'[ \t]+', ' ', text)  # Multiple spaces/tabs to single space
        
        # Remove control characters except newlines and tabs
        text = ''.join(char for char in text if char.isprintable() or char in '\n\t')
        
        return text
    
    def _calculate_content_hash(self, documents: List[Document]) -> str:
        """Calculate hash of document content for deduplication"""
        combined_content = '\n'.join([doc.page_content for doc in documents])
        return hashlib.sha256(combined_content.encode()).hexdigest()
    
    def _get_loader_name(self, file_format: str) -> str:
        """Get the loader name used for a format"""
        loader_map = {
            '.pdf': 'PyPDFLoader',
            '.docx': 'UnstructuredWordDocumentLoader',
            '.doc': 'UnstructuredWordDocumentLoader',
            '.xlsx': 'UnstructuredExcelLoader',
            '.xls': 'UnstructuredExcelLoader',
            '.pptx': 'UnstructuredPowerPointLoader',
            '.ppt': 'UnstructuredPowerPointLoader',
            '.txt': 'TextLoader',
            '.md': 'TextLoader',
            '.html': 'UnstructuredHTMLLoader',
            '.htm': 'UnstructuredHTMLLoader',
            '.csv': 'CSVLoader',
            '.json': 'JSONLoader'
        }
        return loader_map.get(file_format, 'UnknownLoader')
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return self.supported_formats.copy()
    
    def get_loader_stats(self) -> Dict[str, Any]:
        """Get loader statistics"""
        return {
            'supported_formats': len(self.supported_formats),
            'max_file_size_mb': self.max_file_size_mb,
            'ocr_available': TESSERACT_AVAILABLE,
            'office_available': OFFICE_AVAILABLE,
            'formats': self.supported_formats
        }